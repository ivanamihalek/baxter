#! /usr/bin/env python

""" Here we'll import the parts of Django we need. It's recommended to leave
these settings as is, and skip to START OF APPLICATION section below """

# Turn off bytecode generation
import sys
sys.dont_write_bytecode = True

# Django specific settings - needed only in scripts that use django-orm
import os
os.environ.setdefault('DJANGO_SETTINGS_MODULE', 'settings')
import django
django.setup()

import os
from datetime import datetime

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN

from models.bad_bac_models import AntibioticResMutation, PDBStructure, Decoy
from models.bad_bac_models import Gene2UCSCAssembly
import random
import shutil


def title_slide(prs, title):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    subtitle = slide.placeholders[1]

    subtitle.text = datetime.today().strftime('%B %d, %Y') + "\nStudenti:"


def title_and_basic_text(slide, title, basic_text):
    slide.shapes.title.text = title
    # slide.shapes.title.font.size =  Pt(16)
    # Change the font size of the title
    for paragraph in slide.shapes.title.text_frame.paragraphs:
        for run in paragraph.runs:  # Iterate through runs in each paragraph
            run.font.size = Pt(24)  # Change this value to adjust font size

    left_regular = Inches(0.3)  # Position from left
    top_regular = Inches(1.5)  # Position from top
    width_regular = Inches(8)  # Width of the text box
    height_regular = Inches(1)  # Height of the text box

    regular_text_box = slide.shapes.add_textbox(left_regular, top_regular, width_regular, height_regular)
    regular_text_frame = regular_text_box.text_frame
    regular_text_frame.auto_size = None  # Prevent auto-sizing
    regular_text_frame.text = basic_text
    regular_text_frame.word_wrap = True  # Enable word wrap
    # Set alignment for regular text to left
    for paragraph in regular_text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT


def bulleted_list(slide, bullet_points, font_size):
    # Add another text box for bulleted list
    left_bullet = Inches(0.8)  # Position from left
    top_bullet  = Inches(2.0)  # Position below the first text box
    width_bullet = Inches(9)  # Width of the text box
    height_bullet = Inches(3)  # Height of the text box

    bullet_text_box = slide.shapes.add_textbox(left_bullet, top_bullet, width_bullet, height_bullet)
    bullet_text_frame = bullet_text_box.text_frame
    bullet_text_frame.auto_size = None
    bullet_text_frame.word_wrap = True  # Enable word wrap

    # Clear existing paragraphs if needed (optional)
    bullet_text_frame.clear()  # Keeps one empty paragraph

    # Add bullet points
    for i, bullet_point in enumerate(bullet_points):
        p = bullet_text_frame.add_paragraph()
        p.text = f"{i+1}: " + bullet_point
        p.space_after = Inches(0.1)  # Optional: add space after each bullet point
        p.font.size = Pt(font_size)
        p.alignment = PP_ALIGN.LEFT


def species_slide(prs, fingerprint, decoy1, decoy2):

    title_only_layout = prs.slide_layouts[5]  # 1 is the index for title and content layout
    slide = prs.slides.add_slide(title_only_layout)
    title_and_basic_text(slide, "Odredjivanje vrsta genetskom analizom",
                         "Koristeci Flongle u terenskom radu, detektirali smo tri sekvence:")
    font_size = 8
    bullet_points = [fingerprint, decoy1, decoy2]
    random.shuffle(bullet_points)
    bulleted_list(slide, bullet_points, font_size)


def blast_slide_first(prs):
    title_only_layout = prs.slide_layouts[5]  # 1 is the index for title and content layout
    slide = prs.slides.add_slide(title_only_layout)
    title = "Odredjivanje vrsta genetskom analizom"
    basic_text  = "Za svaki od tri slijeda napravili smo pretragu na NCBI BLAST stranici."
    basic_text += "(https://blast.ncbi.nlm.nih.gov/Blast.cgi)"
    title_and_basic_text(slide, title, basic_text)


def blast_slide_screenshot(prs, seq_num):
    title_only_layout = prs.slide_layouts[5]  # 1 is the index for title and content layout
    slide = prs.slides.add_slide(title_only_layout)
    title = "Odredjivanje vrsta genetskom analizom"
    basic_text  = f"Rezultati blast pretrage za slijed {seq_num} (slika zaslona [screenshot])"
    title_and_basic_text(slide, title, basic_text)


def blast_slide_last(prs):
    title_only_layout = prs.slide_layouts[5]  # 1 is the index for title and content layout
    slide = prs.slides.add_slide(title_only_layout)
    title = "Odredjivanje vrsta genetskom analizom"
    basic_text  = "... te zakljucili da bi se u uzorku mogla nalaziti slijedeca patogena vrsta"
    title_and_basic_text(slide, title, basic_text)

    font_size = 16
    bullet_points = ["objasnjenje1", "objasnjenje2", "objasnjenje3"]
    bulleted_list(slide, bullet_points, font_size)


def alignment_slide(prs, refseq_assembly_id):
    title_only_layout = prs.slide_layouts[5]  # 1 is the index for title and content layout
    slide = prs.slides.add_slide(title_only_layout)
    title = "Poravnanje sekvenci"
    basic_text  = "Odnijeli smo uzorak u labos, te izolirali i sekvencirali patogen. "
    basic_text += f"U labosu smo preciznije otkrili da patogen odgovara soju s referentnim sklopom {refseq_assembly_id}. "
    basic_text += "Rezultate sekvenciranja (sample_reads_R1.fastq and sample_reads_R2.fastq) "
    basic_text += "ucitali smo u Galaxy server, https://usegalaxy.org/ (upute: https://shorturl.at/Zfzr5), te napravili poravnanje koristeci BWA program."
    basic_text += " [screenshot]."
    title_and_basic_text(slide, title, basic_text)


def variant_calling_slide(prs, contig, sequenced_range):
    title_only_layout = prs.slide_layouts[5]  # 1 is the index for title and content layout
    slide = prs.slides.add_slide(title_only_layout)
    title = "Pozivanje varijanata"
    basic_text  = "Poravnate readove ucitali smo u Interactive Genomic Viewer (https://igv.org/app/)"
    basic_text  += "da bismo otkrili pozicije varijanata iz naseg uzorka."
    range_def = f"{contig}:{sequenced_range[0]}-{sequenced_range[1]}"
    basic_text  += f"Iz labosa su nam javili da je posebno temeljito sekvencirano podrucje {range_def} "
    basic_text  += f"na genomu patogena. "
    basic_text  += f"Upute za IGV: https://shorturl.at/XorRU. [Screenshot]"
    title_and_basic_text(slide, title, basic_text)


def variants_slide(prs):
    title_only_layout = prs.slide_layouts[5]  # 1 is the index for title and content layout
    slide = prs.slides.add_slide(title_only_layout)
    title = "Pozivanje varijanata"
    basic_text  = "Koristeci IGV, prepoznali smo slijedece varijante:"
    title_and_basic_text(slide, title, basic_text)

    font_size = 16
    bullet_points = ["varijanta 1 ", "varijanta 2", "varijanta3"]
    bulleted_list(slide, bullet_points, font_size)


def ucsc_slide(prs):
    title_only_layout = prs.slide_layouts[5]  # 1 is the index for title and content layout
    slide = prs.slides.add_slide(title_only_layout)
    title = "Interpretacija varijanata"
    basic_text  = "Koji je utjecaj varijante ne genetskom nivou? Utjece li varijanta na kodirajuci gen? "
    basic_text  += "Koji gen? Na koji nacin? \n"
    basic_text  += "https://hgdownload.soe.ucsc.edu/hubs/bacteria/index.html [screenshot]"
    title_and_basic_text(slide, title, basic_text)


def pubmed_slide(prs):
    title_only_layout = prs.slide_layouts[5]  # 1 is the index for title and content layout
    slide = prs.slides.add_slide(title_only_layout)
    title = "Postojece znanje - varijante u literaturi"
    basic_text  = "Pretrazili smo pubmed  https://pubmed.ncbi.nlm.nih.gov/ koristeci ime patogene vrste, "
    basic_text += "gene i proteinsku mutaciju kao upit, te saznali da ta mutacija rezultira u ... [screenshot]"
    title_and_basic_text(slide, title, basic_text)


def pdb_slide(prs):
    title_only_layout = prs.slide_layouts[5]  # 1 is the index for title and content layout
    slide = prs.slides.add_slide(title_only_layout)
    title = "Posljedica varijante "
    basic_text  = "Upisali smo ime gena, patogene vrste te antibiotika u trazilicu "
    basic_text += "PDB-a, repozitorija proteinskih strukura, https://www.rcsb.org/, "
    basic_text += "te pronasli protein kristaliziran s tom molekulom u strukturi [screenshot]"
    title_and_basic_text(slide, title, basic_text)


def mutation_impact_slide(prs):
    title_only_layout = prs.slide_layouts[5]  # 1 is the index for title and content layout
    slide = prs.slides.add_slide(title_only_layout)
    title = "Posljedica varijante "
    basic_text  = "Pregledom strukture (upute: https://shorturl.at/uGEKV) i lokacije mutacije na proteinu,  "
    basic_text += "ustanovili smo da je njen utjecaj ... [screenshot]"
    title_and_basic_text(slide, title, basic_text)


def story_to_pptx(arm_entry, pdb_entries, exercise_label, outdir):
    pptx_filepath = f"{outdir}/vjezba_{exercise_label}.pptx"
    pdb_ids = [pe.pdb_id for pe in pdb_entries]

    aa_from = arm_entry.mutation[0]
    aa_to = arm_entry.mutation[-1]
    pos = int(arm_entry.mutation[1:-1])
    # dist = pdb_2_arm_entry.dist_to_drug

    print(f"\n******************** EXERCISE {exercise_label}  *******************************")
    assmb_entry = arm_entry.assemblies.first()
    print(arm_entry.gene.name, arm_entry.mutation, assmb_entry.refseq_assembly_id, pdb_ids)
    for drug in arm_entry.drugs_affected.all():
        print(drug.name)
    gene_entry = arm_entry.gene
    gene_2_assmb_entry: Gene2UCSCAssembly = Gene2UCSCAssembly.objects.get(gene_id=gene_entry.id, assembly_id=assmb_entry.id)
    contig = gene_2_assmb_entry.contig
    sequenced_range = (gene_2_assmb_entry.start_on_contig-int(2.e4), gene_2_assmb_entry.end_on_contig+int(2.e4))
    prs = Presentation()
    title_slide(prs, f"Vježba {exercise_label}")

    ########################################
    # FINDING THE SPECIES
    # select region from  assmb_entry.refseq_assembly_id, + 2 decoys  --> task: find problematic species using NCBI Blast
    fingerprint = random.choice(assmb_entry.fingerprint_set.all()).dna_seq
    decoy1 = random.choice(Decoy.objects.all()).dna_seq
    decoy2 = random.choice(Decoy.objects.all()).dna_seq
    # q = f"""Q1:  In the field, using your Flongle, you detected the following three DNA snippets:
    # {fingerprint[:10]}   {decoy1[:10]}   {decoy2[:10]}
    # Which one is problematic?  Which species does it belong to? """
    # a = f"{assmb_entry.common_name}  "
    # print(q, "\n", a)
    species_slide(prs, fingerprint, decoy1, decoy2)
    blast_slide_first(prs)
    for i in range(3):
        blast_slide_screenshot(prs, i+1)
    blast_slide_last(prs)

    ########################################
    # DETECTING THE VARIANT
    # toy output from sequencer - the name xxx  --> task: find variants using Galaxy
    q = f"""Q2: You took the sample to the lab for the full sequencing. The lab returned the following sequencing file. 
    # The lab also determined that the strain that this sample contains is {assmb_entry.refseq_assembly_id}
    #  Which variants does this strain cary? What are their genomic coords?"""
    # a = f" genomic coords"
    # print(q, "\n", a)
    alignment_slide(prs, assmb_entry.refseq_assembly_id)
    variant_calling_slide(prs, contig,  sequenced_range)
    variants_slide(prs)
    ########################################
    # PLACING THE VARIANT IN GENOMIC CONTEXT
    # ---> task: find genomic location from the previous step using UCSC Genome Browser and some calculation - gene and protein location
    # q = """Q3: Which genomic locations do these variants hit? Which gene, if any?
    # Which codon is hit on the protein level, if any?"""
    # a = f"{arm_entry.gene.name}, {aa_from} at the position {pos}"
    # print(q, "\n", a)

    ########################################
    # DETERMINING THE VARIANT CONSEQUENCE ON GENOMIC LEVEL
    # ---> task what is the impact of mutation - use codon translation table Wikipedia
    # q = """Q4: What is the impact of the nucleotide change? """
    # a = f"{aa_from}>{aa_to}"
    # print(q, "\n", a)
    ucsc_slide(prs)

    ########################################
    # SEARCHING FOR THE PREVIOUS KNOWLEDGE
    # ---> is this mutation known?  CARD database -- should get the drug name
    # q = """Q5: Is this mutation already known in the literature? Which drug resistance it may cause?"""
    # drugs = [drug.name for drug in arm_entry.drugs_affected.all()]
    # a = f"{drugs} pubmed ids -- any should do"
    # print(q, "\n", a)
    pubmed_slide(prs)

    ########################################
    # UNDERSTANDING THE IMPACT ON THE SYSTEMC LEVEL
    # ---> PDB search for species + gene + drug name -- make illustration
    # q = """Q6: Is the structure of the protein-drug complex known for this bacterial species?
    # Use it to inspect the impact that the mutation will have on the drug binding."""
    # a = f"any of the pdbs that I have here, {pdb_ids}"
    # print(q, "\n", a)
    pdb_slide(prs)
    mutation_impact_slide(prs)

    prs.save(pptx_filepath)
    return pptx_filepath


def run():
    scratch_dir = "/home/ivana/scratch/baxter"
    bac_genomes = "/storage/databases/ucsc/bacterial_genomes"
    # "arm" = antibiotic resistance mutation

    # start with pdb, arm pairs where the mutation is close to ligand
    # this is supposed to save us time from investigating arms that

    ecount = 0
    # for pdb_2_arm_entry in Pdb2Mutation.objects.filter(dist_to_drug__lte=9.0):
    for arm_entry in AntibioticResMutation.objects.filter(assemblies__isnull=False).distinct():
        pdb_entries = PDBStructure.objects.filter(
            pdb2mutation__antibio_res_mutation=arm_entry,
            pdb2mutation__dist_to_drug__lt=9
        ).distinct()
        if len(pdb_entries) < 1: continue

        outdir = f"{scratch_dir}/toy_genomes/{arm_entry.gene.name}_{arm_entry.mutation}"
        if not os.path.exists(outdir): continue
        ecount += 1
        pptx_filepath = story_to_pptx(arm_entry, pdb_entries, ecount, outdir)
        print(pptx_filepath)

        # copy reference genome fa & fai to otdir
        assmb_entry = arm_entry.assemblies.first()

        shutil.copy(f"{bac_genomes}/{assmb_entry.refseq_assembly_id}.fa", outdir)
        shutil.copy(f"{bac_genomes}/{assmb_entry.refseq_assembly_id}.fa.fai", outdir)

        # store to db?
        # question_and_answers
        # category, title, question, answers, specific


#######################
if __name__ == "__main__":
    run()
